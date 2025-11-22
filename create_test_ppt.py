from pptx import Presentation
from pptx.util import Inches

# Create a new presentation
prs = Presentation()

# Add a title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Test Presentation"
subtitle.text = "This is a test slide for PPT to PDF conversion"

# Add a second slide with content
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)

title = slide.shapes.title
title.text = "Second Slide"

body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "This is the main point"

p = tf.add_paragraph()
p.text = "This is a second level bullet"
p.level = 1

p = tf.add_paragraph()
p.text = "This is a third level bullet"
p.level = 2

# Add an image (optional, just to make it more complex)
# We'll just keep it simple for testing

# Save the presentation
prs.save('test_presentation.pptx')
print("Created test_presentation.pptx")