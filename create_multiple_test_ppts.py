from pptx import Presentation
from pptx.util import Inches

def create_test_presentation(filename, title):
    """Create a simple test presentation with the given filename and title"""
    prs = Presentation()

    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle = slide.placeholders[1]

    title_shape.text = title
    subtitle.text = f"This is the test slide for {filename}"

    # Add a second slide with content
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = "Second Slide"

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "This is the main point"

    p = tf.add_paragraph()
    p.text = "This is a second level bullet"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "This is a third level bullet"
    p.level = 2

    # Save the presentation
    prs.save(filename)
    print(f"Created {filename}")

# Create multiple test files for bulk conversion
test_files = [
    ("presentation_1.pptx", "First Test Presentation"),
    ("presentation_2.pptx", "Second Test Presentation"),
    ("presentation_3.pptx", "Third Test Presentation"),
    ("my_presentation.pptx", "My Sample Presentation")
]

for filename, title in test_files:
    create_test_presentation(filename, title)

print("\nAll test presentations created successfully!")